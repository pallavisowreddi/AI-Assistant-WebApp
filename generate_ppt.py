import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide (Layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Nexus AI"
    slide1.placeholders[1].text = (
        "Prompt Engineering & Feedback Assistant\n\n"
        "VaultOfCodes Internship Major Project\n"
        "Submitted by: Pallavi Sowreddi\n\n"
        "Live Website: ai-assistant-web-app-nine.vercel.app\n"
        "GitHub: github.com/pallavisowreddi/AI-Assistant-WebApp"
    )
    
    # Slide 2: Problem Statement (Layout 1)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Problem Statement"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Developing optimal, structured prompts for LLMs is difficult for general users."
    
    p = tf2.add_paragraph()
    p.text = "• Traditional chat interfaces do not allow users to easily toggle and test prompt variants (e.g. Factual vs. Pedagogical)."
    p.level = 0
    
    p = tf2.add_paragraph()
    p.text = "• Developers and analysts lack integrated metrics to evaluate which system prompts yield the highest user satisfaction."
    p.level = 0
    
    p = tf2.add_paragraph()
    p.text = "• Absence of real-time 'Prompt Inspectors' to educate students on how prompt wrappers shape LLM outputs under the hood."
    p.level = 0

    # Slide 3: Proposed Solution (Layout 1)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Proposed Solution"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Built a web-based playground mapping user inputs to engineered system templates."
    
    p = tf3.add_paragraph()
    p.text = "• Developed 12 custom prompt blueprints across 4 distinct cognitive workloads."
    p.level = 0
    
    p = tf3.add_paragraph()
    p.text = "• Integrated a visual 'Prompt Inspector' showing the final payload sent to the API."
    p.level = 0
    
    p = tf3.add_paragraph()
    p.text = "• Implemented a user feedback database (JSON) to track satisfaction ratings and comments."
    p.level = 0
    
    # Slide 4: Core Features (Layout 1)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Project Features"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "A full-featured dashboard optimized for prompt tuning:"
    
    p = tf4.add_paragraph()
    p.text = "• 4 Modules: Question Answering, Text Summarization, Content Generation, Study Advisor."
    p.level = 0
    
    p = tf4.add_paragraph()
    p.text = "• Interactive Controls: Prompt variants toggle, rating clickers, and text reviews."
    p.level = 0
    
    p = tf4.add_paragraph()
    p.text = "• Analytics Dashboard: Live totals, Helpful Rate bar charts, and scrollable log logs."
    p.level = 0
    
    p = tf4.add_paragraph()
    p.text = "• Export Utilities: Copy Prompt and Export Logs to Excel-compatible CSV formats."
    p.level = 0

    # Slide 5: Technology & Tools Used (Layout 1)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Technology Stack"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Modern full-stack technologies with zero external hosting dependencies:"
    
    p = tf5.add_paragraph()
    p.text = "• Core Backend: Python & Flask framework."
    p.level = 0
    
    p = tf5.add_paragraph()
    p.text = "• Generative AI API: Google Gemini API (gemini-2.5-flash model)."
    p.level = 0
    
    p = tf5.add_paragraph()
    p.text = "• Frontend Design: Responsive HTML5, Vanilla CSS3 (Neon-accented Glassmorphism), JavaScript (AJAX)."
    p.level = 0
    
    p = tf5.add_paragraph()
    p.text = "• Deployment & Versioning: Git, GitHub, Vercel."
    p.level = 0

    # Slide 6: Screenshots & Demo Links (Layout 1)
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Project Deployment & Code Links"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Access our live deployment and open-source code:"
    
    p = tf6.add_paragraph()
    p.text = "• Live Deployed Website Link: https://ai-assistant-web-app-nine.vercel.app/"
    p.level = 0
    
    p = tf6.add_paragraph()
    p.text = "• GitHub Repository URL: https://github.com/pallavisowreddi/AI-Assistant-WebApp"
    p.level = 0
    
    p = tf6.add_paragraph()
    p.text = "[PASTE YOUR SCREENSHOTS / DEMO IMAGES HERE IN POWERPOINT]"
    p.level = 0
    
    # Slide 7: Conclusion (Layout 1)
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Conclusion"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "This project succeeds in delivering a fully functional Prompt Engineering benchmark:"
    
    p = tf7.add_paragraph()
    p.text = "• Showcases how system constraints shape model responses (length, tone, output structures)."
    p.level = 0
    
    p = tf7.add_paragraph()
    p.text = "• Bridges prompt testing with quantitative user feedback and data exports."
    p.level = 0
    
    p = tf7.add_paragraph()
    p.text = "• Ready-to-evaluate, production-grade codebase deployed for public review."
    p.level = 0

    # Save presentation
    output_filename = "AI_Assistant_Project_Presentation.pptx"
    prs.save(output_filename)
    print(f"PowerPoint created successfully: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_presentation()
